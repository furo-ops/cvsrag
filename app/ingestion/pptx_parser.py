import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

logger = logging.getLogger(__name__)

# MSO_SHAPE_TYPE.TABLE = 19
_TABLE_SHAPE_TYPE = 19


def extract_text_from_pptx(file_path: str) -> dict:
    """Extract all text from a .pptx file, returns structured dict."""
    path = Path(file_path)

    try:
        prs = Presentation(str(path))
    except Exception as e:
        logger.error(f"Cannot open {path.name}: {e}")
        raise

    slides_content = []
    all_text_parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []

        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)

                if shape.shape_type == _TABLE_SHAPE_TYPE:
                    table = shape.table
                    for row in table.rows:
                        row_texts = [
                            cell.text.strip()
                            for cell in row.cells
                            if cell.text.strip()
                        ]
                        if row_texts:
                            slide_texts.append(" | ".join(row_texts))
            except Exception as e:
                logger.warning(f"Error reading shape in slide {slide_num} of {path.name}: {e}")
                continue

        if slide_texts:
            slides_content.append({"slide": slide_num, "text": "\n".join(slide_texts)})
            all_text_parts.extend(slide_texts)

    # Infer name from filename: "john_doe_cv.pptx" -> "John Doe"
    stem = path.stem
    for suffix in ("_cv", "-cv", "_CV", "-CV"):
        stem = stem.replace(suffix, "")
    name = stem.replace("_", " ").replace("-", " ").title()

    return {
        "name": name,
        "source_file": path.name,
        "raw_text": "\n".join(all_text_parts),
        "slides_content": slides_content,
    }
