#!/usr/bin/env python3
"""
Generate 10 fake consulting-profile .pptx CVs and a matching availability.csv.

Usage:
    python scripts/generate_sample_data.py
"""

import csv
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT_DIR = Path("./data/cvs")
AVAILABILITY_FILE = Path("./data/availability.csv")

GREEN = RGBColor(0x86, 0xBC, 0x25)
NAVY = RGBColor(0x1A, 0x1A, 0x2E)

PROFILES = [
    {
        "filename": "ana_garcia_cv.pptx",
        "name": "Ana García",
        "title": "Senior Data Engineer",
        "summary": (
            "Senior Data Engineer with 6 years of experience designing and building "
            "large-scale data pipelines on Azure. Specialist in Databricks, Apache Spark, "
            "and real-time streaming architectures. Strong Python and SQL background."
        ),
        "skills": [
            "Python", "Apache Spark", "Databricks", "Azure Data Factory",
            "Azure Synapse Analytics", "SQL", "Delta Lake", "Kafka", "dbt",
            "Azure DevOps", "Terraform", "Power BI",
        ],
        "certifications": [
            "Microsoft Certified: Azure Data Engineer Associate (DP-203)",
            "Databricks Certified Associate Developer for Apache Spark",
        ],
        "education": "BSc Computer Science — Universidad Complutense de Madrid",
        "languages": ["Python", "SQL", "Spanish (native)", "English (C1)"],
        "experience": [
            ("2021–Present", "Senior Data Engineer", "Deloitte", "Led migration of on-prem data warehouse to Azure Synapse for a major retail client."),
            ("2019–2021", "Data Engineer", "Accenture", "Built real-time Kafka pipelines processing 10M events/day for a telco client."),
            ("2018–2019", "Junior Developer", "Indra", "Developed ETL processes for financial reporting systems."),
        ],
    },
    {
        "filename": "carlos_martinez_cv.pptx",
        "name": "Carlos Martínez",
        "title": "Cloud Solutions Architect",
        "summary": (
            "Cloud Architect with 9 years of experience delivering enterprise cloud solutions "
            "on Azure and AWS. Expert in multi-cloud architecture, infrastructure-as-code, "
            "microservices, and cloud cost optimization. Azure certified at expert level."
        ),
        "skills": [
            "Azure", "AWS", "Terraform", "Kubernetes", "Docker", "Azure AKS",
            "Azure API Management", "Microservices", "Bicep", "CI/CD",
            "GitHub Actions", "Python", "Bash",
        ],
        "certifications": [
            "Microsoft Certified: Azure Solutions Architect Expert (AZ-305)",
            "AWS Certified Solutions Architect – Professional",
            "Microsoft Certified: DevOps Engineer Expert (AZ-400)",
        ],
        "education": "MSc Telecommunications Engineering — Universidad Politécnica de Madrid",
        "languages": ["Python", "Bash", "Spanish (native)", "English (C2)", "French (B1)"],
        "experience": [
            ("2020–Present", "Cloud Architect", "Deloitte", "Designed multi-cloud architecture for global logistics company reducing infra costs by 35%."),
            ("2017–2020", "Cloud Engineer", "Capgemini", "Led Azure migrations for 5 enterprise clients in financial services sector."),
            ("2015–2017", "Systems Engineer", "Telefónica", "Managed on-prem infrastructure for core business applications."),
        ],
    },
    {
        "filename": "maria_lopez_cv.pptx",
        "name": "María López",
        "title": "AI/ML Specialist",
        "summary": (
            "AI and Machine Learning specialist with 5 years of experience building and "
            "deploying production ML systems. Deep expertise in NLP, LLMs, and MLOps on Azure. "
            "Published researcher with background in neural networks and time-series forecasting."
        ),
        "skills": [
            "Python", "PyTorch", "TensorFlow", "scikit-learn", "Azure Machine Learning",
            "MLflow", "LangChain", "OpenAI API", "Hugging Face", "NLP",
            "Azure OpenAI Service", "pandas", "FastAPI", "Docker",
        ],
        "certifications": [
            "Microsoft Certified: Azure AI Engineer Associate (AI-102)",
            "Microsoft Certified: Azure Data Scientist Associate (DP-100)",
            "Deep Learning Specialization — Coursera/DeepLearning.AI",
        ],
        "education": "MSc Artificial Intelligence — Universidad Autónoma de Madrid",
        "languages": ["Python", "R", "Spanish (native)", "English (C1)"],
        "experience": [
            ("2022–Present", "AI/ML Specialist", "Deloitte AI Institute", "Developed RAG-based internal knowledge assistant using Azure OpenAI and Cognitive Search."),
            ("2020–2022", "Data Scientist", "BBVA", "Built credit risk scoring models using gradient boosting and neural networks."),
            ("2019–2020", "ML Engineer (Intern)", "Telefónica I+D", "Research on LSTM models for network traffic prediction."),
        ],
    },
    {
        "filename": "juan_sanchez_cv.pptx",
        "name": "Juan Sánchez",
        "title": "Data Scientist",
        "summary": (
            "Data Scientist with 4 years of experience in advanced analytics, statistical modeling, "
            "and machine learning. Specialist in Power BI storytelling and Python-based ML pipelines. "
            "Strong business acumen — translates complex models into executive insights."
        ),
        "skills": [
            "Python", "R", "scikit-learn", "Power BI", "DAX", "SQL",
            "Azure Machine Learning", "pandas", "matplotlib", "seaborn",
            "Statistical modeling", "A/B testing", "Jupyter",
        ],
        "certifications": [
            "Microsoft Certified: Power BI Data Analyst Associate (PL-300)",
            "Google Data Analytics Professional Certificate",
        ],
        "education": "BSc Statistics — Universidad Carlos III de Madrid",
        "languages": ["Python", "R", "SQL", "Spanish (native)", "English (B2)"],
        "experience": [
            ("2022–Present", "Data Scientist", "Deloitte Analytics", "Delivered churn prediction model for telecom client, reducing churn by 18%."),
            ("2020–2022", "Business Analyst", "Santander", "Built executive dashboards and KPI reporting systems in Power BI."),
        ],
    },
    {
        "filename": "laura_fernandez_cv.pptx",
        "name": "Laura Fernández",
        "title": "Senior Project Manager / Scrum Master",
        "summary": (
            "Certified PMP and Scrum Master with 8 years managing technology transformation projects "
            "in financial services, retail, and public sector. Expert in hybrid Agile methodologies, "
            "risk management, and stakeholder communication across international teams."
        ),
        "skills": [
            "Project Management", "Agile", "Scrum", "SAFe", "Risk Management",
            "Jira", "Confluence", "MS Project", "Stakeholder Management",
            "Change Management", "Azure DevOps", "Kanban",
        ],
        "certifications": [
            "Project Management Professional (PMP)",
            "Certified Scrum Master (CSM)",
            "SAFe 6 Practitioner",
        ],
        "education": "MBA — IE Business School, Madrid",
        "languages": ["Spanish (native)", "English (C2)", "Portuguese (B2)"],
        "experience": [
            ("2019–Present", "Senior PM / Scrum Master", "Deloitte", "Programme lead for €5M cloud transformation at major Spanish bank, delivered on time and budget."),
            ("2016–2019", "Project Manager", "EY", "Managed regulatory compliance projects (GDPR, PSD2) for financial sector clients."),
            ("2014–2016", "Business Analyst", "Indra", "Requirements analysis and PMO support for public sector IT projects."),
        ],
    },
    {
        "filename": "david_gonzalez_cv.pptx",
        "name": "David González",
        "title": "Power BI Developer / Data Analyst",
        "summary": (
            "Power BI specialist with 5 years of experience designing enterprise BI solutions. "
            "Expert in data modeling, DAX, and Power Query. Experienced with Azure Synapse "
            "as data backend and embedding Power BI in custom applications."
        ),
        "skills": [
            "Power BI", "DAX", "Power Query", "SQL", "Azure Synapse Analytics",
            "Azure Analysis Services", "Excel", "Python", "Data Modeling",
            "Tableau", "Paginated Reports", "SSRS",
        ],
        "certifications": [
            "Microsoft Certified: Power BI Data Analyst Associate (PL-300)",
            "Microsoft Certified: Azure Data Fundamentals (DP-900)",
        ],
        "education": "BSc Business Administration — Universidad de Sevilla",
        "languages": ["SQL", "DAX", "Spanish (native)", "English (B2)"],
        "experience": [
            ("2021–Present", "Power BI Developer", "Deloitte", "Built enterprise BI platform for retail chain with 200+ reports and 500 daily users."),
            ("2019–2021", "Data Analyst", "Inditex", "Designed sales analytics dashboards for global operations team."),
        ],
    },
    {
        "filename": "elena_romero_cv.pptx",
        "name": "Elena Romero",
        "title": "MLOps / DevOps Engineer",
        "summary": (
            "MLOps and DevOps engineer with 6 years specializing in ML model deployment, "
            "monitoring, and CI/CD for data science teams. Expert in Azure ML pipelines, "
            "Kubernetes, and building reproducible ML infrastructure at scale."
        ),
        "skills": [
            "MLflow", "Azure Machine Learning", "Kubeflow", "Kubernetes",
            "Docker", "Terraform", "GitHub Actions", "Azure DevOps",
            "Python", "Bash", "Prometheus", "Grafana", "Helm",
        ],
        "certifications": [
            "Microsoft Certified: Azure DevOps Engineer Expert (AZ-400)",
            "Certified Kubernetes Administrator (CKA)",
            "Microsoft Certified: Azure AI Engineer Associate (AI-102)",
        ],
        "education": "BSc Software Engineering — Universidad de Valencia",
        "languages": ["Python", "Bash", "Spanish (native)", "English (C1)"],
        "experience": [
            ("2021–Present", "MLOps Engineer", "Deloitte AI", "Built ML platform on AKS serving 15 production models with automated retraining pipelines."),
            ("2019–2021", "DevOps Engineer", "NTT Data", "Implemented GitOps workflows and container orchestration for microservices migration."),
            ("2018–2019", "Systems Admin", "GMV", "Linux infrastructure management and automation with Ansible."),
        ],
    },
    {
        "filename": "miguel_torres_cv.pptx",
        "name": "Miguel Torres",
        "title": "SAP Senior Consultant",
        "summary": (
            "SAP consultant with 10 years of experience implementing SAP ERP, S/4HANA, and "
            "BTP solutions for manufacturing and retail companies. Specialist in SAP FI/CO, MM, "
            "and integration with cloud platforms via SAP Integration Suite."
        ),
        "skills": [
            "SAP S/4HANA", "SAP ERP", "SAP FI/CO", "SAP MM", "SAP BTP",
            "SAP Integration Suite", "ABAP", "SAP Fiori", "SAP Analytics Cloud",
            "ITSM", "ITIL", "SQL",
        ],
        "certifications": [
            "SAP Certified Application Associate – SAP S/4HANA Finance",
            "SAP Certified Technology Associate – SAP Integration Suite",
            "ITIL 4 Foundation",
        ],
        "education": "BSc Industrial Engineering — Universidad de Zaragoza",
        "languages": ["ABAP", "SQL", "Spanish (native)", "English (B2)", "German (A2)"],
        "experience": [
            ("2018–Present", "SAP Senior Consultant", "Deloitte", "Led S/4HANA greenfield implementation for €800M manufacturing company, 18-month project."),
            ("2015–2018", "SAP Consultant", "Accenture", "SAP ERP rollouts for retail clients across Spain and Portugal."),
            ("2014–2015", "Junior SAP Analyst", "Seidor", "SAP FI/CO configuration and end-user training."),
        ],
    },
    {
        "filename": "sofia_diaz_cv.pptx",
        "name": "Sofía Díaz",
        "title": "ServiceNow Developer",
        "summary": (
            "ServiceNow certified developer with 5 years implementing ITSM, ITOM, and custom "
            "workflows on the ServiceNow platform. Experience integrating ServiceNow with Azure "
            "AD, Teams, and third-party APIs. Strong JavaScript and REST API skills."
        ),
        "skills": [
            "ServiceNow", "ITSM", "ITOM", "JavaScript", "REST API",
            "Azure Active Directory", "Microsoft Teams", "Glide", "Flow Designer",
            "Integration Hub", "CMDB", "Python",
        ],
        "certifications": [
            "ServiceNow Certified Application Developer",
            "ServiceNow Certified Implementation Specialist – ITSM",
            "Microsoft Certified: Azure Fundamentals (AZ-900)",
        ],
        "education": "BSc Information Systems — Universidad de Málaga",
        "languages": ["JavaScript", "Python", "Spanish (native)", "English (C1)"],
        "experience": [
            ("2022–Present", "ServiceNow Developer", "Deloitte", "Designed custom ITSM portal for 5,000-user enterprise, integrating with Azure AD SSO."),
            ("2020–2022", "ServiceNow Admin", "DXC Technology", "Platform administration and workflow automation for financial services client."),
        ],
    },
    {
        "filename": "pablo_moreno_cv.pptx",
        "name": "Pablo Moreno",
        "title": "Cybersecurity Consultant",
        "summary": (
            "Cybersecurity consultant with 7 years in cloud security, penetration testing, "
            "and security architecture. Specialist in Azure security controls, Zero Trust, "
            "and Microsoft Sentinel SIEM. Experienced with ISO 27001 and ENS compliance frameworks."
        ),
        "skills": [
            "Azure Security", "Microsoft Sentinel", "Zero Trust", "Penetration Testing",
            "SIEM", "Azure Defender for Cloud", "ISO 27001", "ENS",
            "Entra ID", "Conditional Access", "Threat Intelligence", "Python",
        ],
        "certifications": [
            "Microsoft Certified: Security Operations Analyst Associate (SC-200)",
            "Microsoft Certified: Azure Security Engineer Associate (AZ-500)",
            "Certified Ethical Hacker (CEH)",
            "OSCP (Offensive Security Certified Professional)",
        ],
        "education": "MSc Cybersecurity — Universidad Rey Juan Carlos",
        "languages": ["Python", "Bash", "Spanish (native)", "English (C1)"],
        "experience": [
            ("2020–Present", "Cybersecurity Consultant", "Deloitte Cyber", "Led security architecture review and Zero Trust roadmap for national bank."),
            ("2018–2020", "Security Analyst", "S21Sec", "SOC L2 analyst, incident response, and threat hunting."),
            ("2017–2018", "IT Security Intern", "Bankinter", "Vulnerability management and security awareness training."),
        ],
    },
]

AVAILABILITY_DATA = [
    {"name": "Ana García",     "current_project": "Digital Transformation – Retail Bank",  "availability_date": "2026-04-01", "availability_percentage": 0,   "location": "Madrid",    "grade": "Senior Consultant"},
    {"name": "Carlos Martínez","current_project": "",                                         "availability_date": "2026-03-03", "availability_percentage": 100, "location": "Barcelona", "grade": "Manager"},
    {"name": "María López",    "current_project": "AI Assistant – Insurance",                "availability_date": "2026-03-17", "availability_percentage": 30,  "location": "Madrid",    "grade": "Senior Consultant"},
    {"name": "Juan Sánchez",   "current_project": "",                                         "availability_date": "2026-03-03", "availability_percentage": 100, "location": "Sevilla",   "grade": "Consultant"},
    {"name": "Laura Fernández","current_project": "ERP Migration – Manufacturing",            "availability_date": "2026-06-01", "availability_percentage": 0,   "location": "Madrid",    "grade": "Senior Manager"},
    {"name": "David González", "current_project": "BI Platform – Retail",                    "availability_date": "2026-03-31", "availability_percentage": 20,  "location": "Madrid",    "grade": "Consultant"},
    {"name": "Elena Romero",   "current_project": "",                                         "availability_date": "2026-03-03", "availability_percentage": 100, "location": "Valencia",  "grade": "Senior Consultant"},
    {"name": "Miguel Torres",  "current_project": "S/4HANA Rollout – Manufacturing",         "availability_date": "2026-05-15", "availability_percentage": 0,   "location": "Zaragoza",  "grade": "Senior Consultant"},
    {"name": "Sofía Díaz",     "current_project": "ServiceNow ITSM – Utilities",             "availability_date": "2026-04-15", "availability_percentage": 0,   "location": "Málaga",    "grade": "Consultant"},
    {"name": "Pablo Moreno",   "current_project": "",                                         "availability_date": "2026-03-03", "availability_percentage": 100, "location": "Madrid",    "grade": "Manager"},
]


def add_slide(prs: Presentation, title: str, content_lines: list[str]) -> None:
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.color.rgb = NAVY
    title_shape.text_frame.paragraphs[0].font.size = Pt(20)
    title_shape.text_frame.paragraphs[0].font.bold = True

    # Content
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)


def create_cv(profile: dict, output_dir: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = profile["name"]
    slide.placeholders[1].text = profile["title"]
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = NAVY
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True

    # Summary
    add_slide(prs, "Professional Summary", [profile["summary"]])

    # Skills
    add_slide(prs, "Technical Skills", profile["skills"])

    # Certifications
    add_slide(prs, "Certifications", profile["certifications"])

    # Education & Languages
    add_slide(prs, "Education & Languages", [profile["education"]] + profile["languages"])

    # Experience
    exp_lines = []
    for years, role, company, desc in profile["experience"]:
        exp_lines.append(f"{years} | {role} @ {company}")
        exp_lines.append(f"  {desc}")
        exp_lines.append("")
    add_slide(prs, "Professional Experience", exp_lines)

    out_path = output_dir / profile["filename"]
    prs.save(str(out_path))
    print(f"  Created: {out_path}")


def create_availability_csv(data: list[dict], path: Path) -> None:
    with open(path, "w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(
            f,
            fieldnames=["name", "current_project", "availability_date",
                         "availability_percentage", "location", "grade"],
        )
        writer.writeheader()
        writer.writerows(data)
    print(f"  Created: {path}")


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    AVAILABILITY_FILE.parent.mkdir(parents=True, exist_ok=True)

    print("Generating sample CV files...")
    for profile in PROFILES:
        create_cv(profile, OUTPUT_DIR)

    print("\nGenerating availability.csv...")
    create_availability_csv(AVAILABILITY_DATA, AVAILABILITY_FILE)

    print(f"\nDone. {len(PROFILES)} CVs + availability.csv generated.")
    print("Run: python scripts/ingest_cvs.py")


if __name__ == "__main__":
    main()
